#!/usr/bin/env python3
"""Buat PPT bimbingan TA — progress SOAR open-source.
Versi diperbaiki: bahasa lebih natural, tidak terlalu formal/robotik."""
import sys
sys.path.insert(0, "/tmp/pptx-env/lib/python3.13/site-packages")

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─── Warna (lebih soft, tidak terlalu neon) ───
BG = RGBColor(0xF5, 0xF5, 0xF0)
DARK = RGBColor(0x2D, 0x2D, 0x2D)
BLUE = RGBColor(0x1A, 0x56, 0x8E)
GREEN_DARK = RGBColor(0x2E, 0x7D, 0x32)
ORANGE = RGBColor(0xE6, 0x5C, 0x00)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xEE, 0xEE, 0xEE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED_DARK = RGBColor(0xC6, 0x28, 0x28)


def add_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox


def add_bullet_list(slide, left, top, width, height, items, size=16, color=DARK, spacing=6):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_before = Pt(spacing)
    return txBox


def add_card(slide, left, top, width, height, title, items, title_color=BLUE, bg_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = LIGHT_GRAY
    shape.line.width = Pt(1)
    shape.shadow.inherit = False

    txBox = slide.shapes.add_textbox(Inches(left + 0.25), Inches(top + 0.15), Inches(width - 0.5), Inches(height - 0.3))
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.color.rgb = title_color
    p.font.bold = True
    p.space_after = Pt(6)

    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(13)
        p.font.color.rgb = DARK
        p.space_before = Pt(3)


def add_divider(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Pt(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()


# ═══════════════════════════════════════════════════════════════
# SLIDE 1: Cover
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, RGBColor(0x1A, 0x56, 0x8E))

add_text(slide, 1, 1.8, 11, 1, "Sistem SOAR Open-Source", 44, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, 1, 2.8, 11, 1, "Berbasis n8n untuk Deteksi & Respons Ancaman\nMalware dan Phishing", 24, RGBColor(0xBB, 0xDE, 0xFB), False, PP_ALIGN.CENTER)

add_text(slide, 1, 4.8, 11, 0.5, "───────────────────────────", 16, RGBColor(0x90, 0xCA, 0xF9), False, PP_ALIGN.CENTER)

add_text(slide, 1, 5.3, 11, 0.5, "Ravi Arnan Irianto  |  2305551076", 20, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, 1, 5.9, 11, 0.5, "Bimbingan Tugas Akhir  ·  September 2026", 16, RGBColor(0xBB, 0xDE, 0xFB), False, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2: Apa yang sudah dikerjakan
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, 0.6, 0.3, 12, 0.7, "Apa yang Sudah Dikerjakan", 30, BLUE, True)
add_divider(slide, 0.6, 0.95, 5)

# Kolom kiri: bug fix + detection
add_card(slide, 0.5, 1.3, 5.8, 2.6, "Bug Fix & Deteksi", [
    "block-domain sekarang persist (tidak hilang saat restart container)",
    "Notifikasi ganda untuk 1 file sudah diperbaiki",
    "File berekstensi berisiko (.sh/.exe/.ps1) yang tak dikenal VT -> otomatis minta review",
    "File tanpa ekstensi tapi punya execute bit juga tertangkap",
], BLUE)

# Kolom kanan: threat intel
add_card(slide, 6.8, 1.3, 5.8, 2.6, "Threat Intelligence", [
    "Sekarang pakai 2 sumber: VirusTotal + MalwareBazaar",
    "Cache VT pakai TTL diferensial (bersih 24 jam, berbahaya 7 hari)",
    "Phishing proaktif: fetch URLhaus tiap jam, block sebelum user klik",
    "Kalau VT error/rate-limit, sistem tetap jalan dan kasih tahu analis",
], BLUE)

# Baris bawah: infra + benchmark
add_card(slide, 0.5, 4.2, 5.8, 2.6, "Infrastruktur & Deployment", [
    "n8n di-update ke 2.36.9 (di atas semua CVE 2026)",
    "Hardening: Caddy reverse-proxy + TLS + basic auth",
    "IaC pake Ansible playbook (idempoten)",
    "Health monitor: cek agent/n8n/Ollama, alert cuma saat status berubah",
], GREEN_DARK)

add_card(slide, 6.8, 4.2, 5.8, 2.6, "Pengukuran & Evaluasi", [
    "Benchmark udah jalan: throughput 34 alert/detik",
    "FN rate 0% (15/15 file berisiko ketangkep semua)",
    "FP suppression 100% (8 alert baseline -> 0 notifikasi)",
    "Mapping ke MITRE ATT&CK: 10 teknik tercakup",
], GREEN_DARK)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3: Cara Kerja Sistem
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, 0.6, 0.3, 12, 0.7, "Cara Kerja Sistem", 30, BLUE, True)
add_divider(slide, 0.6, 0.95, 4)

# Flow sederhana
flow = [
    ("Endpoint\n(FIM/Log)", 0.3),
    ("Wazuh\nManager", 2.6),
    ("n8n\nWorkflow", 4.9),
    ("VirusTotal\n+ MalwareBazaar", 7.2),
    ("Ollama\n(AI lokal)", 9.8),
    ("Telegram\n(analis)", 12.1),
]
for label, x in flow:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.4), Inches(2.0), Inches(1.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

# panah
for i in range(len(flow) - 1):
    x1 = flow[i][1] + 2.0
    x2 = flow[i+1][1]
    add_text(slide, x1, 1.7, x2 - x1, 0.4, "->", 20, BLUE, True, PP_ALIGN.CENTER)

# Decision table
add_text(slide, 0.6, 2.8, 12, 0.5, "Keputusan berdasarkan keyakinan:", 18, BLUE, True)

table_data = [
    ("Situasi", "Yang dilakukan"),
    ("VT deteksi >= 20", "Langsung isolasi file otomatis"),
    ("VT deteksi 5-19 atau MB match", "Kirim ke Telegram, analis pilih tombol"),
    ("VT deteksi 1-4 atau file risky unknown", "Kirim + minta saran AI"),
    ("VT error / rate-limit", "Tandai degradasi, minta verifikasi manual"),
    ("VT bersih + MB bersih", "Diam saja (tidak ganggu analis)"),
]

table = slide.shapes.add_table(len(table_data), 2, Inches(0.5), Inches(3.3), Inches(12), Inches(3.2)).table
table.columns[0].width = Inches(5)
table.columns[1].width = Inches(7)

for row_idx, (col1, col2) in enumerate(table_data):
    for col_idx, text in enumerate([col1, col2]):
        cell = table.cell(row_idx, col_idx)
        cell.text = text
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.color.rgb = DARK
            if row_idx == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
        if row_idx == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE
        elif row_idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE


# ═══════════════════════════════════════════════════════════════
# SLIDE 4: Hasil Pengukuran
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, 0.6, 0.3, 12, 0.7, "Hasil Pengukuran", 30, BLUE, True)
add_divider(slide, 0.6, 0.95, 4)

add_text(slide, 0.6, 1.2, 12, 0.5, "Uji coba dijalankan langsung di sistem live (bukan simulasi)", 15, GRAY)

# MTTR
add_card(slide, 0.5, 1.8, 4.0, 2.3, "Kecepatan Respons (MTTR)", [
    "Malware (N=30): 1,68 detik",
    "Phishing (N=10): 2,13 detik",
    "Webhook response: 0,03 detik",
    "Median malware: 1,42 detik",
], BLUE)

# Throughput
add_card(slide, 4.7, 1.8, 4.0, 2.3, "Kapasitas (Load Test)", [
    "34 alert per detik (N=20)",
    "Rata-rata 142 ms per alert",
    "5 request paralel",
    "Total 0,6 detik untuk 20 alert",
], BLUE)

# FN Rate
add_card(slide, 8.9, 1.8, 4.0, 2.3, "Akurasi Deteksi", [
    "FN rate: 0% (N=15)",
    "File risky + hash unknown = 15/15 ketangkep",
    "FP suppression: 100% (N=8)",
    "Baseline 8 alert -> 0 notifikasi",
], GREEN_DARK)

# Konteks
add_card(slide, 0.5, 4.4, 12.3, 2.5, "Konteks Pengukuran", [
    "MTTR 1,68 detik itu waktu dari alert masuk sampai file terisolasi (VT cache hangat).",
    "Kalau VT belum pernah lihat hash itu (cold), butuh ~3-5 detik tambahan untuk scan.",
    "Phishing lebih lama sedikit karena ada pengecekan URLScan (tapi GSB langsung).",
    "Load test 34 alert/detik itu jauh di atas beban normal SOC (biasanya 1-5 alert/menit).",
    "FN rate 0% artinya semua file berisiko yang kami uji berhasil terdeteksi, tidak ada yang lolos.",
], GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5: Kontribusi / Kebaruan
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, 0.6, 0.3, 12, 0.7, "Kontribusi & Kebaruan", 30, BLUE, True)
add_divider(slide, 0.6, 0.95, 4)

add_card(slide, 0.5, 1.3, 5.8, 2.5, "Kenapa ini berbeda dari SOAR lain?", [
    "Sebagian besar SOAR (Shuffle, Cortex) itu black-box — analis tidak tahu kenapa keputusan diambil.",
    "Sistem ini kasih alasan di setiap notifikasi: kenapa file ini dianggap berbahaya.",
    "Kalau sumber intel error, sistem jujur bilang 'saya tidak bisa verifikasi' bukan diam saja.",
    "Keputusan auto-isolate cuma untuk keyakinan tinggi, sisanya serahin ke analis.",
], BLUE)

add_card(slide, 6.8, 1.3, 5.8, 2.5, "Yang belum ada di SOAR lain", [
    "LLM advisory lokal (Ollama) untuk alert yang ambigu — tanpa kirim data ke cloud.",
    "Phishing proaktif: blok URL dari feed publik sebelum user sempat klik.",
    "Health monitor yang sadar diri: tahu kapan ia sedang tidak bisa bekerja dengan baik.",
    "Cache VT dengan TTL diferensial: file bersih di-scan ulang lebih cepat.",
], BLUE)

add_card(slide, 0.5, 4.1, 5.8, 2.8, "Bukti pendukung", [
    "Mapping MITRE ATT&CK: 10 teknik tercakup dalam playbook.",
    "Perbandingan empiris n8n vs Shuffle: 6 aspek dinilai (n8n menang 3,8 vs 2,5).",
    "Benchmark N>=30 dengan data nyata, bukan simulasi.",
    "Semua konfigurasi di-version-control (Ansible + Docker Compose).",
], GREEN_DARK)

add_card(slide, 6.8, 4.1, 5.8, 2.8, "Fitur yang paling berguna di dunia nyata", [
    "Self-aware: kalau VT lagi down, notifikasi langsung kasih tahu analis.",
    "Audit trail: semua keputusan analis tercatat (siapa, kapan, apa keputusannya).",
    "Phishing proaktif: seringkali URL sudah diblokir sebelum user buka email.",
    "FP suppression 100%: analis tidak dibanjiri alert palsu.",
], GREEN_DARK)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6: Rencana ke Depan
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, 0.6, 0.3, 12, 0.7, "Rencana ke Depan", 30, BLUE, True)
add_divider(slide, 0.6, 0.95, 4)

items = [
    ("Trusted Autonomy", "Sekarang: semua keputusan butuh konfirmasi analis (HITL).\nRencana: kalau keyakinan sangat tinggi, boleh auto tanpa tunggu analis, tapi ada timeout & SLA.", ORANGE),
    ("RAG (Retrieval-Augmented Generation)", "Masalah: LLM kadang mengarang rekomendasi (halusinasi).\nRencana: kasih LLM akses ke playbook & threat-intel supaya jawabannya berbasis data.", ORANGE),
    ("Arsitektur HA", "Sekarang: 1 container n8n + SQLite (single point of failure).\nRencana: n8n queue-mode + Redis + PostgreSQL + Prometheus/Grafana.", RED_DARK),
    ("Upgrade Wazuh", "4.9.2 -> 4.14.7. Deferred pasca-TA karena agent juga harus di-upgrade bareng.", GRAY),
]

for i, (title, desc, color) in enumerate(items):
    y = 1.2 + i * 1.5
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(12.3), Inches(1.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = color
    shape.line.width = Pt(2)

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(y + 0.1), Inches(11.7), Inches(1.1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(17)
    p.font.color.rgb = color
    p.font.bold = True

    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(13)
    p2.font.color.rgb = DARK
    p2.space_before = Pt(4)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7: Penutup
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, RGBColor(0x1A, 0x56, 0x8E))

add_text(slide, 1, 2.2, 11, 1, "Terima Kasih", 44, WHITE, True, PP_ALIGN.CENTER)

add_text(slide, 1.5, 3.8, 10, 1.2,
    "Sistem ini dibangun dengan prinsip:\n"
    "bukan menggantikan analis, tapi membantu analis\n"
    "ambil keputusan lebih cepat dengan informasi yang cukup.",
    18, RGBColor(0xBB, 0xDE, 0xFB), False, PP_ALIGN.CENTER)

add_text(slide, 1, 5.5, 11, 0.5, "Ravi Arnan Irianto  |  2305551076", 18, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, 1, 6.0, 11, 0.5, "raviarnankeren@gmail.com", 14, RGBColor(0x90, 0xCA, 0xF9), False, PP_ALIGN.CENTER)


# ─── Simpan ───
output_path = "docs/BIMBINGAN-TA-SOAR.pptx"
prs.save(output_path)
print(f"PPT disimpan ke {output_path}")
print(f"Total slides: {len(prs.slides)}")
